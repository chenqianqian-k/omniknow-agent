import hashlib
from pathlib import Path

from docx import Document as DocxDocument
from pptx import Presentation

from langchain_chroma import Chroma
from langchain_community.document_loaders import PyPDFLoader
from langchain_core.documents import Document as LangChainDocument
from langchain_huggingface import HuggingFaceEmbeddings
from langchain_text_splitters import RecursiveCharacterTextSplitter


# =====================================================
# 项目路径
# =====================================================

PROJECT_ROOT = Path(__file__).resolve().parents[2]

UPLOAD_DIR = PROJECT_ROOT / "data" / "uploads"
CHROMA_DIR = PROJECT_ROOT / "data" / "chroma"

EMBEDDING_MODEL = (
    PROJECT_ROOT
    / "models"
    / "bge-small-zh-v1.5"
)

COLLECTION_NAME = "knowledge_base"


# =====================================================
# 支持的文档类型
# =====================================================

SUPPORTED_EXTENSIONS = {
    ".pdf",
    ".docx",
    ".pptx",
    ".md",
    ".markdown",
    ".txt",
}


# =====================================================
# 通用辅助函数
# =====================================================

def calculate_file_hash(
    file_path: str,
) -> str:
    """计算文件的SHA-256哈希值。"""
    hash_calculator = hashlib.sha256()

    with open(file_path, "rb") as file:
        while chunk := file.read(8192):
            hash_calculator.update(chunk)

    return hash_calculator.hexdigest()


def get_file_type(
    path: Path,
) -> str:
    """根据文件扩展名返回文档类型。"""
    extension = path.suffix.lower()

    if extension == ".markdown":
        return "md"

    return extension.lstrip(".")


def get_base_metadata(
    path: Path,
    file_hash: str,
) -> dict:
    """创建不同文档格式共用的元数据。"""
    return {
        "file_name": path.name,
        "file_type": get_file_type(path),
        "source": str(path),
        "file_hash": file_hash,
    }


def get_embeddings() -> HuggingFaceEmbeddings:
    """加载文本向量模型。"""
    return HuggingFaceEmbeddings(
        model_name=str(EMBEDDING_MODEL),
        model_kwargs={
            "device": "cpu",
        },
        encode_kwargs={
            "normalize_embeddings": True,
        },
    )


def get_vector_store() -> Chroma:
    """连接持久化的ChromaDB向量库。"""
    return Chroma(
        collection_name=COLLECTION_NAME,
        embedding_function=get_embeddings(),
        persist_directory=str(CHROMA_DIR),
    )


# =====================================================
# PDF解析
# =====================================================

def load_pdf(
    path: Path,
    file_hash: str,
) -> list[LangChainDocument]:
    """解析PDF，每一页转换为一个Document。"""
    loader = PyPDFLoader(str(path))
    documents = loader.load()

    base_metadata = get_base_metadata(
        path=path,
        file_hash=file_hash,
    )

    for index, document in enumerate(documents):
        page_index = document.metadata.get(
            "page",
            index,
        )

        document.metadata.update(
            base_metadata
        )

        # 保留原来的零起始page字段，
        # 兼容当前retriever.py中的page + 1逻辑。
        document.metadata["page"] = int(
            page_index
        )

        document.metadata[
            "location_type"
        ] = "page"

        document.metadata[
            "location"
        ] = int(page_index) + 1

        document.metadata[
            "location_label"
        ] = f"第{int(page_index) + 1}页"

    return documents


# =====================================================
# Word解析
# =====================================================

def load_docx(
    path: Path,
    file_hash: str,
) -> list[LangChainDocument]:
    """解析DOCX中的段落和表格文字。"""
    docx_document = DocxDocument(str(path))

    content_parts = []

    # 读取Word中的普通段落。
    for paragraph in docx_document.paragraphs:
        text = paragraph.text.strip()

        if text:
            content_parts.append(text)

    # 读取Word表格中的文字。
    for table_index, table in enumerate(
        docx_document.tables,
        start=1,
    ):
        table_rows = []

        for row in table.rows:
            cell_texts = [
                cell.text.strip()
                for cell in row.cells
            ]

            if any(cell_texts):
                table_rows.append(
                    " | ".join(cell_texts)
                )

        if table_rows:
            content_parts.append(
                f"\n【表格{table_index}】\n"
                + "\n".join(table_rows)
            )

    content = "\n\n".join(content_parts)

    if not content.strip():
        raise ValueError(
            f"Word文档中没有提取到文字：{path.name}"
        )

    metadata = get_base_metadata(
        path=path,
        file_hash=file_hash,
    )

    metadata.update(
        {
            # 暂时保留page字段，
            # 保证当前retriever.py能够正常工作。
            "page": 0,
            "location_type": "document",
            "location": 1,
            "location_label": "文档正文",
        }
    )

    return [
        LangChainDocument(
            page_content=content,
            metadata=metadata,
        )
    ]


# =====================================================
# PowerPoint解析
# =====================================================

def load_pptx(
    path: Path,
    file_hash: str,
) -> list[LangChainDocument]:
    """解析PPTX，每张幻灯片转换为一个Document。"""
    presentation = Presentation(str(path))

    documents = []

    for slide_index, slide in enumerate(
        presentation.slides,
        start=1,
    ):
        slide_content = []

        for shape in slide.shapes:
            # 读取普通文本框、标题和形状中的文字。
            if (
                hasattr(shape, "has_text_frame")
                and shape.has_text_frame
            ):
                text = shape.text.strip()

                if text:
                    slide_content.append(text)

            # 读取幻灯片表格中的文字。
            if (
                hasattr(shape, "has_table")
                and shape.has_table
            ):
                table_rows = []

                for row in shape.table.rows:
                    cell_texts = [
                        cell.text.strip()
                        for cell in row.cells
                    ]

                    if any(cell_texts):
                        table_rows.append(
                            " | ".join(cell_texts)
                        )

                if table_rows:
                    slide_content.append(
                        "\n".join(table_rows)
                    )

        content = "\n\n".join(slide_content)

        # 没有文字的幻灯片暂时跳过。
        if not content.strip():
            continue

        metadata = get_base_metadata(
            path=path,
            file_hash=file_hash,
        )

        metadata.update(
            {
                # 使用零起始值兼容当前retriever。
                "page": slide_index - 1,
                "location_type": "slide",
                "location": slide_index,
                "location_label": (
                    f"第{slide_index}张幻灯片"
                ),
            }
        )

        documents.append(
            LangChainDocument(
                page_content=content,
                metadata=metadata,
            )
        )

    if not documents:
        raise ValueError(
            f"PPTX中没有提取到文字：{path.name}"
        )

    return documents


# =====================================================
# Markdown和TXT解析
# =====================================================

def read_text_file(
    path: Path,
) -> str:
    """尝试使用常见编码读取文本文件。"""
    encodings = [
        "utf-8-sig",
        "utf-8",
        "gb18030",
    ]

    last_error = None

    for encoding in encodings:
        try:
            return path.read_text(
                encoding=encoding
            )

        except UnicodeDecodeError as error:
            last_error = error

    raise ValueError(
        f"无法识别文件编码：{path.name}"
    ) from last_error


def load_text_document(
    path: Path,
    file_hash: str,
) -> list[LangChainDocument]:
    """解析Markdown或TXT文本文件。"""
    content = read_text_file(path).strip()

    if not content:
        raise ValueError(
            f"文本文件内容为空：{path.name}"
        )

    metadata = get_base_metadata(
        path=path,
        file_hash=file_hash,
    )

    metadata.update(
        {
            "page": 0,
            "location_type": "document",
            "location": 1,
            "location_label": "文档正文",
        }
    )

    return [
        LangChainDocument(
            page_content=content,
            metadata=metadata,
        )
    ]


# =====================================================
# 统一文档解析入口
# =====================================================

def load_document(
    file_path: str,
    file_hash: str,
) -> list[LangChainDocument]:
    """根据文件扩展名选择对应的文档解析器。"""
    path = Path(file_path)
    extension = path.suffix.lower()

    if extension == ".pdf":
        return load_pdf(
            path=path,
            file_hash=file_hash,
        )

    if extension == ".docx":
        return load_docx(
            path=path,
            file_hash=file_hash,
        )

    if extension == ".pptx":
        return load_pptx(
            path=path,
            file_hash=file_hash,
        )

    if extension in {
        ".md",
        ".markdown",
        ".txt",
    }:
        return load_text_document(
            path=path,
            file_hash=file_hash,
        )

    raise ValueError(
        f"不支持的文件格式：{extension}"
    )


def load_and_split_document(
    file_path: str,
    file_hash: str | None = None,
):
    """解析常见文档格式，并切分为适合检索的文本块。"""
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(
            f"文件不存在：{path}"
        )

    extension = path.suffix.lower()

    if extension not in SUPPORTED_EXTENSIONS:
        supported_types = ", ".join(
            sorted(SUPPORTED_EXTENSIONS)
        )

        raise ValueError(
            "当前不支持该文件格式。"
            f"支持的格式：{supported_types}"
        )

    if file_hash is None:
        file_hash = calculate_file_hash(
            str(path)
        )

    documents = load_document(
        file_path=str(path),
        file_hash=file_hash,
    )

    splitter = RecursiveCharacterTextSplitter(
        chunk_size=600,
        chunk_overlap=100,
        separators=[
            "\n\n",
            "\n",
            "。",
            "！",
            "？",
            ".",
            " ",
            "",
        ],
    )

    chunks = splitter.split_documents(
        documents
    )

    if not chunks:
        raise ValueError(
            f"文档切分后没有可入库内容：{path.name}"
        )

    return documents, chunks


# =====================================================
# 兼容原有PDF函数
# =====================================================

def load_and_split_pdf(
    pdf_path: str,
    file_hash: str | None = None,
):
    """兼容旧代码：解析并切分PDF。"""
    path = Path(pdf_path)

    if path.suffix.lower() != ".pdf":
        raise ValueError(
            "load_and_split_pdf只支持PDF文件"
        )

    return load_and_split_document(
        file_path=pdf_path,
        file_hash=file_hash,
    )


# =====================================================
# 文档入库
# =====================================================

def ingest_document(
    file_path: str,
) -> dict:
    """将文档切分、向量化并写入ChromaDB。"""
    path = Path(file_path)

    if not path.exists():
        raise FileNotFoundError(
            f"文件不存在：{path}"
        )

    extension = path.suffix.lower()

    if extension not in SUPPORTED_EXTENSIONS:
        supported_types = ", ".join(
            sorted(SUPPORTED_EXTENSIONS)
        )

        raise ValueError(
            "当前不支持该文件格式。"
            f"支持的格式：{supported_types}"
        )

    file_hash = calculate_file_hash(
        str(path)
    )

    vector_store = get_vector_store()

    existing_data = vector_store.get(
        where={
            "file_hash": file_hash,
        },
        include=[
            "metadatas",
        ],
    )

    if existing_data.get("ids"):
        existing_metadatas = (
            existing_data.get(
                "metadatas",
                [],
            )
        )

        locations = {
            metadata.get(
                "location_label",
                "文档正文",
            )
            for metadata in existing_metadatas
            if metadata
        }

        first_metadata = (
            existing_metadatas[0]
            if existing_metadatas
            else {}
        )

        unit_count = len(locations) or 1

        return {
            "file_name": path.name,
            "file_type": (
                first_metadata.get(
                    "file_type"
                )
                or get_file_type(path)
            ),
            # 暂时保留page_count，
            # 兼容当前API和frontend.py。
            "page_count": unit_count,
            "unit_count": unit_count,
            "chunk_count": len(
                existing_data["ids"]
            ),
            "collection": COLLECTION_NAME,
            "file_hash": file_hash,
            "already_exists": True,
        }

    documents, chunks = (
        load_and_split_document(
            file_path=str(path),
            file_hash=file_hash,
        )
    )

    ids = [
        f"{file_hash}-{index}"
        for index, _ in enumerate(chunks)
    ]

    vector_store.add_documents(
        documents=chunks,
        ids=ids,
    )

    unit_count = len(documents)

    return {
        "file_name": path.name,
        "file_type": get_file_type(path),
        # 暂时保留page_count，
        # 兼容当前API和frontend.py。
        "page_count": unit_count,
        "unit_count": unit_count,
        "chunk_count": len(chunks),
        "collection": COLLECTION_NAME,
        "file_hash": file_hash,
        "already_exists": False,
    }


def ingest_pdf(
    pdf_path: str,
) -> dict:
    """兼容旧代码：将PDF写入知识库。"""
    path = Path(pdf_path)

    if path.suffix.lower() != ".pdf":
        raise ValueError(
            "ingest_pdf只支持PDF文件"
        )

    return ingest_document(pdf_path)


# =====================================================
# 文档列表
# =====================================================

def list_documents() -> list[dict]:
    """查询知识库中已经入库的文档。"""
    vector_store = get_vector_store()

    stored_data = vector_store.get(
        include=[
            "metadatas",
        ],
    )

    documents = {}

    for metadata in stored_data.get(
        "metadatas",
        [],
    ):
        if not metadata:
            continue

        file_name = metadata.get(
            "file_name",
            "未知文件",
        )

        file_hash = metadata.get(
            "file_hash"
        )

        file_type = metadata.get(
            "file_type",
            Path(file_name).suffix.lower().lstrip("."),
        )

        document_key = (
            file_hash
            or file_name
        )

        if document_key not in documents:
            documents[document_key] = {
                "file_name": file_name,
                "file_type": file_type,
                "file_hash": file_hash,
                "chunk_count": 0,
            }

        documents[
            document_key
        ]["chunk_count"] += 1

    return list(documents.values())


# =====================================================
# 文档删除
# =====================================================

def delete_document(
    file_hash: str | None = None,
    file_name: str | None = None,
) -> dict:
    """删除指定文档在ChromaDB中的全部文本块。"""
    if not file_hash and not file_name:
        raise ValueError(
            "file_hash和file_name至少需要提供一个"
        )

    vector_store = get_vector_store()

    if file_hash:
        search_condition = {
            "file_hash": file_hash,
        }

    else:
        search_condition = {
            "file_name": file_name,
        }

    stored_data = vector_store.get(
        where=search_condition,
        include=[
            "metadatas",
        ],
    )

    chunk_ids = stored_data.get(
        "ids",
        [],
    )

    if not chunk_ids:
        return {
            "deleted": False,
            "file_name": file_name,
            "deleted_chunk_count": 0,
            "message": "没有找到需要删除的文档",
        }

    metadatas = stored_data.get(
        "metadatas",
        [],
    )

    if not file_name and metadatas:
        file_name = metadatas[0].get(
            "file_name",
            "未知文件",
        )

    vector_store.delete(
        ids=chunk_ids
    )

    return {
        "deleted": True,
        "file_name": file_name,
        "deleted_chunk_count": len(
            chunk_ids
        ),
        "message": "文档已从知识库删除",
    }


# =====================================================
# 命令行测试
# =====================================================

if __name__ == "__main__":
    import argparse

    parser = argparse.ArgumentParser(
        description=(
            "将PDF、Word、PPTX、Markdown"
            "或TXT文档写入知识库"
        )
    )

    parser.add_argument(
        "file_path",
        help="需要写入知识库的文档路径",
    )

    args = parser.parse_args()

    result = ingest_document(
        args.file_path
    )

    if result["already_exists"]:
        print("文档已经存在，无需重复入库")

    else:
        print("文档入库成功")

    print(
        f"文件：{result['file_name']}"
    )
    print(
        f"类型：{result['file_type']}"
    )
    print(
        f"内容单元数：{result['unit_count']}"
    )
    print(
        f"文本块数：{result['chunk_count']}"
    )
    print(
        f"集合：{result['collection']}"
    )
